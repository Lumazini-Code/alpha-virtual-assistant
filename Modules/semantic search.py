import json
import numpy as np
import faiss
from pathlib import Path
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation
import socket
import onnxruntime as ort
from tokenizers import Tokenizer

session = ort.InferenceSession(
    "semanticSearch_onnx/model.onnx",
    providers=["CPUExecutionProvider"]
)

tokenizer = Tokenizer.from_file("semanticSearch_onnx/tokenizer.json")
tokenizer.enable_padding(pad_token="[PAD]", pad_id=0, length=128)
tokenizer.enable_truncation(max_length=128)


def encode(texts: list[str]) -> np.ndarray:
    if isinstance(texts, str):
        texts = [texts]

    encoded = tokenizer.encode_batch(texts)

    input_ids      = np.array([e.ids            for e in encoded], dtype=np.int64)
    attention_mask = np.array([e.attention_mask  for e in encoded], dtype=np.int64)

    outputs = session.run(None, {
        "input_ids":      input_ids,
        "attention_mask": attention_mask,
    })

    token_embeddings = outputs[0]
    mask = attention_mask[..., np.newaxis].astype(np.float32)
    embeddings = (token_embeddings * mask).sum(axis=1) / mask.sum(axis=1)
    return embeddings.astype(np.float32)


def startServer():
    HOST = '127.0.0.1'
    PORT = 7
    server_socket = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    # FIX BUG #3 (parcial) + SMELL: evita "Address already in use" ao reiniciar
    server_socket.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
    server_socket.bind((HOST, PORT))
    server_socket.listen(1)
    return server_socket


def aguardar_conexao(server_socket):
    print("[INFO] Aguardando conexão...")
    conn, addr = server_socket.accept()
    print(f"[INFO] Cliente conectado: {addr}")
    return conn


# ── leitores por formato ─────────────────────────────────────────────────────

def ler_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")

def ler_json(path: Path) -> str:
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, ensure_ascii=False, indent=2)
    except Exception:
        return path.read_text(encoding="utf-8", errors="ignore")

def ler_html(path: Path) -> str:
    soup = BeautifulSoup(path.read_bytes(), "html.parser")
    return soup.get_text(separator=" ", strip=True)

def ler_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    return " ".join(page.extract_text() or "" for page in reader.pages)

def ler_docx(path: Path) -> str:
    doc = Document(str(path))
    return " ".join(p.text for p in doc.paragraphs)

def ler_xlsx(path: Path) -> str:
    wb = openpyxl.load_workbook(str(path), data_only=True)
    linhas = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            linha = " ".join(str(c) for c in row if c is not None)
            if linha.strip():
                linhas.append(linha)
    return "\n".join(linhas)

def ler_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    textos = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                textos.append(shape.text)
    return " ".join(textos)


LEITORES = {
    ".txt": ler_txt, ".md": ler_txt, ".mdx": ler_txt, ".rst": ler_txt,
    ".tex": ler_txt, ".asciidoc": ler_txt, ".adoc": ler_txt,
    ".csv": ler_txt, ".tsv": ler_txt, ".json": ler_json, ".jsonl": ler_txt,
    ".xml": ler_txt, ".yaml": ler_txt, ".yml": ler_txt, ".toml": ler_txt,
    ".ini": ler_txt, ".cfg": ler_txt, ".conf": ler_txt, ".env": ler_txt,
    ".properties": ler_txt, ".html": ler_html, ".htm": ler_html,
    ".xhtml": ler_html, ".css": ler_txt, ".scss": ler_txt, ".sass": ler_txt,
    ".less": ler_txt, ".svg": ler_txt, ".py": ler_txt, ".js": ler_txt,
    ".ts": ler_txt, ".jsx": ler_txt, ".tsx": ler_txt, ".java": ler_txt,
    ".c": ler_txt, ".h": ler_txt, ".cpp": ler_txt, ".hpp": ler_txt,
    ".cs": ler_txt, ".go": ler_txt, ".rs": ler_txt, ".rb": ler_txt,
    ".php": ler_txt, ".swift": ler_txt, ".kt": ler_txt, ".r": ler_txt,
    ".m": ler_txt, ".lua": ler_txt, ".pl": ler_txt, ".sh": ler_txt,
    ".bash": ler_txt, ".zsh": ler_txt, ".ps1": ler_txt, ".bat": ler_txt,
    ".cmd": ler_txt, ".vbs": ler_txt, ".sql": ler_txt, ".graphql": ler_txt,
    ".proto": ler_txt, ".dart": ler_txt, ".ex": ler_txt, ".exs": ler_txt,
    ".erl": ler_txt, ".hs": ler_txt, ".scala": ler_txt, ".clj": ler_txt,
    ".vim": ler_txt, ".dockerfile": ler_txt, ".gitignore": ler_txt,
    ".editorconfig": ler_txt, ".htaccess": ler_txt, ".nginx": ler_txt,
    ".log": ler_txt, ".out": ler_txt, ".err": ler_txt, ".diff": ler_txt,
    ".patch": ler_txt, ".pdf": ler_pdf, ".docx": ler_docx,
    ".xlsx": ler_xlsx, ".xls": ler_xlsx, ".pptx": ler_pptx,
}


# ── funções principais ───────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 300, overlap: int = 50) -> list[str]:
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk:
            chunks.append(chunk)
    return chunks


def indexar_documentos(arquivos: list, output_dir: str = "index"):
    # FIX BUG #1: converter para Path logo no início
    output_dir = Path(output_dir)

    todos_chunks = []
    metadados = []

    if not arquivos:
        return

    for arquivo in arquivos:
        arquivo = Path(arquivo)
        if arquivo.suffix.lower() not in LEITORES:
            print(f"  [aviso] formato não suportado: {arquivo.name}")
            continue
        print(f"  lendo {arquivo.name} ({arquivo.suffix})")
        leitor = LEITORES[arquivo.suffix.lower()]
        texto = leitor(arquivo)
        if not texto.strip():
            print(f"  [aviso] {arquivo.name} veio vazio, pulando")
            continue
        for chunk in chunk_text(texto):
            todos_chunks.append(chunk)
            metadados.append({"arquivo": str(arquivo), "trecho": chunk})

    if not todos_chunks:
        print("Nenhum conteúdo extraído.")
        return

    print(f"\nIndexando {len(todos_chunks)} trechos de {len(arquivos)} arquivo(s)...")
    vetores = encode(todos_chunks)

    index = faiss.IndexFlatIP(vetores.shape[1])
    faiss.normalize_L2(vetores)
    index.add(vetores)

    output_dir.mkdir(exist_ok=True)
    faiss.write_index(index, str(output_dir / "index.faiss"))
    with open(output_dir / "metadados.json", "w", encoding="utf-8") as f:
        json.dump(metadados, f, ensure_ascii=False, indent=2)

    print("Indexação concluída!")


def carregar_index(output_dir: str = "index"):
    # FIX BUG #1: converter para Path logo no início
    output_dir = Path(output_dir)
    index = faiss.read_index(str(output_dir / "index.faiss"))
    with open(output_dir / "metadados.json", encoding="utf-8") as f:
        metadados = json.load(f)
    return index, metadados


def buscar(pergunta: str, index, metadados, top_k: int = 3) -> list[dict]:
    vetor = encode([pergunta])
    faiss.normalize_L2(vetor)
    scores, indices = index.search(vetor, top_k)
    return [
        {"trecho": metadados[indices[0][i]]["trecho"]}
        for i in range(len(indices[0]))
    ]


# ── receber mensagem completa com delimitador ────────────────────────────────

def recv_mensagem(conn: socket.socket, buffer: list) -> str | None:
    """
    FIX BUG #2 (fragmentação TCP): acumula dados no buffer até encontrar '\n'.
    O C# deve terminar cada mensagem enviada com '\n'.
    Retorna a mensagem sem o '\n', ou None se a conexão fechou.
    """
    while True:
        if b"\n" in buffer[0]:
            idx = buffer[0].index(b"\n")
            msg = buffer[0][:idx].decode("utf-8", errors="ignore").strip()
            buffer[0] = buffer[0][idx + 1:]
            return msg
        chunk = conn.recv(4096)
        if not chunk:
            return None  # conexão fechada
        buffer[0] += chunk


# ── loop principal ───────────────────────────────────────────────────────────

server_socket = startServer()

index = None
metadados = None

while True:
    # FIX BUG #2 (reconexão): aceita nova conexão sempre que o cliente desconecta
    conn = aguardar_conexao(server_socket)
    buffer = [b""]

    while True:
        user_input = recv_mensagem(conn, buffer)

        if user_input is None:
            # cliente desconectou — volta ao accept()
            print("[INFO] Cliente desconectado.")
            conn.close()
            break

        if not user_input:
            continue

        if user_input.startswith("<files>"):
            conteudo = user_input[len("<files>"):].strip().lstrip(";").strip()
            arquivos = [Path(p.strip()) for p in conteudo.split(";") if p.strip()]
            indexar_documentos(arquivos)
            index, metadados = carregar_index()
            conn.sendall("<indexed>\n".encode("utf-8"))

        elif user_input.startswith("<search>"):
            ask = user_input[len("<search>"):].strip().lstrip(";").strip()
            if ask:
                if index is None or metadados is None:
                    conn.sendall("Nenhum arquivo indexado ainda.{end}\n".encode("utf-8"))
                    continue
                resultados = buscar(ask, index, metadados)
                for r in resultados:
                    resposta = f"Trecho: {r['trecho']}\n"   # \n aqui
                    conn.sendall(resposta.encode("utf-8"))
                conn.sendall("{end}\n".encode("utf-8"))